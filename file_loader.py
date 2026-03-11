"""
文件加载模块 - 支持 PDF、Word (.docx)、TXT、PPT (.pptx)、Markdown (.md) 文件解析
统一接口：load_file(uploaded_file) -> list[Document]

PDF 增强：使用 pdfplumber 提取表格，使用 PyMuPDF 提取图片中的文字描述
"""
import os
import tempfile
from langchain_core.documents import Document
from langchain_community.document_loaders import PyMuPDFLoader
from docx import Document as DocxDocument
from pptx import Presentation
import re

def _save_temp_file(uploaded_file, suffix):
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getvalue())
        return tmp.name


def _load_pdf(uploaded_file) -> list[Document]:
    """加载 PDF：PyMuPDF 提取文本 + pdfplumber 提取表格 + PyMuPDF 提取图片信息"""
    tmp_path = _save_temp_file(uploaded_file, ".pdf")
    try:
        # 1. 基础文本提取（PyMuPDF）
        loader = PyMuPDFLoader(tmp_path)
        text_docs = loader.load()

        # 2. 表格提取（pdfplumber）
        table_docs = _extract_tables(tmp_path, uploaded_file.name)

        # 3. 图片描述提取（PyMuPDF）
        image_docs = _extract_image_info(tmp_path, uploaded_file.name)

        return text_docs + table_docs + image_docs
    finally:
        os.remove(tmp_path)


def _extract_tables(pdf_path: str, source_name: str) -> list[Document]:
    """使用 pdfplumber 从 PDF 提取所有表格，转为 Markdown 格式"""
    try:
        import pdfplumber
    except ImportError:
        return []

    documents = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables):
                    if not table or len(table) < 2:
                        continue
                    # 转为 Markdown 表格
                    md_table = _table_to_markdown(table)
                    if md_table:
                        documents.append(Document(
                            page_content=f"[表格 {table_idx + 1}]\n{md_table}",
                            metadata={
                                "source": source_name,
                                "page": page_num,
                                "content_type": "table",
                            }
                        ))
    except Exception:
        pass
    return documents


def _table_to_markdown(table: list[list]) -> str:
    """将二维表格数据转为 Markdown 格式"""
    if not table:
        return ""

    # 清理单元格
    cleaned = []
    for row in table:
        cleaned_row = [(cell or "").replace("\n", " ").strip() for cell in row]
        cleaned.append(cleaned_row)

    if not cleaned:
        return ""

    # 构建 Markdown
    header = cleaned[0]
    col_count = len(header)
    lines = ["| " + " | ".join(header) + " |"]
    lines.append("| " + " | ".join(["---"] * col_count) + " |")
    for row in cleaned[1:]:
        # 补齐列数
        while len(row) < col_count:
            row.append("")
        lines.append("| " + " | ".join(row[:col_count]) + " |")

    return "\n".join(lines)


def _extract_image_info(pdf_path: str, source_name: str) -> list[Document]:
    """从 PDF 中提取图片的基本信息（尺寸、位置、所在页）"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return []

    documents = []
    try:
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images(full=True)
            if not images:
                continue

            # 提取该页图片周围的文字作为上下文
            page_text = page.get_text("text")
            for img_idx, img in enumerate(images):
                xref = img[0]
                img_info = doc.extract_image(xref)
                if img_info:
                    width = img_info.get("width", 0)
                    height = img_info.get("height", 0)
                    ext = img_info.get("ext", "unknown")

                    # 找到图片附近的文字块作为描述
                    caption = _find_image_caption(page, page_text, img_idx)

                    content_parts = [
                        f"[图片 {img_idx + 1}] 第 {page_num + 1} 页",
                        f"尺寸: {width}x{height}, 格式: {ext}",
                    ]
                    if caption:
                        content_parts.append(f"图片描述/标题: {caption}")

                    documents.append(Document(
                        page_content="\n".join(content_parts),
                        metadata={
                            "source": source_name,
                            "page": page_num + 1,
                            "content_type": "image",
                        }
                    ))
        doc.close()
    except Exception:
        pass
    return documents


def _find_image_caption(page, page_text: str, img_idx: int) -> str:
    """尝试找到图片的标题/说明文字（如 Fig. 1, 图 1 等）"""
    # 在页面文字中搜索常见图片标题模式
    patterns = [
        rf'(?:Fig(?:ure)?|图)\s*\.?\s*{img_idx + 1}\s*[.:：]\s*(.+?)(?:\n|$)',
        rf'(?:Fig(?:ure)?|图)\s*\.?\s*{img_idx + 1}\s*(.+?)(?:\n|$)',
    ]
    for pattern in patterns:
        match = re.search(pattern, page_text, re.IGNORECASE)
        if match:
            return match.group(1).strip()[:200]
    return ""


def _load_docx(uploaded_file) -> list[Document]:
    tmp_path = _save_temp_file(uploaded_file, ".docx")
    try:
        doc = DocxDocument(tmp_path)
        documents = []
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if text:
                documents.append(Document(
                    page_content=text,
                    metadata={"source": uploaded_file.name, "paragraph": i + 1}
                ))
        if len(documents) > 20:
            merged = []
            for start in range(0, len(documents), 10):
                batch = documents[start:start + 10]
                content = "\n".join(d.page_content for d in batch)
                merged.append(Document(
                    page_content=content,
                    metadata={"source": uploaded_file.name, "page": start // 10 + 1}
                ))
            return merged
        return documents
    finally:
        os.remove(tmp_path)


def _load_txt(uploaded_file) -> list[Document]:
    """加载 TXT 文件"""
    raw = uploaded_file.getvalue()
    # 尝试多种编码
    for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            text = raw.decode(encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        text = raw.decode("utf-8", errors="ignore")

    return [Document(
        page_content=text,
        metadata={"source": uploaded_file.name, "page": 1}
    )]


def _load_pptx(uploaded_file) -> list[Document]:
    """加载 PPT (.pptx) 文件"""
    tmp_path = _save_temp_file(uploaded_file, ".pptx")
    try:
        prs = Presentation(tmp_path)
        documents = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                documents.append(Document(
                    page_content="\n".join(texts),
                    metadata={"source": uploaded_file.name, "page": slide_num}
                ))
        return documents
    finally:
        os.remove(tmp_path)


def _load_markdown(uploaded_file) -> list[Document]:
    """加载 Markdown (.md) 文件，按标题分段"""
    raw = uploaded_file.getvalue()
    for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            text = raw.decode(encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        text = raw.decode("utf-8", errors="ignore")

    sections = re.split(r'(?=^#{1,3}\s)', text, flags=re.MULTILINE)
    documents = []
    for i, section in enumerate(sections):
        content = section.strip()
        if content:
            documents.append(Document(
                page_content=content,
                metadata={"source": uploaded_file.name, "page": i + 1}
            ))

    if not documents:
        documents.append(Document(
            page_content=text,
            metadata={"source": uploaded_file.name, "page": 1}
        ))
    return documents


_LOADERS = {
    "pdf": _load_pdf,
    "docx": _load_docx,
    "txt": _load_txt,
    "pptx": _load_pptx,
    "md": _load_markdown,
}

SUPPORTED_EXTENSIONS = list(_LOADERS.keys())


def load_file(uploaded_file) -> list[Document]:
    file_name = uploaded_file.name
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""

    loader_fn = _LOADERS.get(ext)
    if loader_fn is None:
        raise ValueError(f"不支持的文件格式: .{ext}（支持: {', '.join(SUPPORTED_EXTENSIONS)}）")

    return loader_fn(uploaded_file)
